"""
make_slides.py — Generate a PowerPoint summary of today's NewtonBench work.
Run: python make_slides.py
Opens the file automatically when done.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Colour palette ─────────────────────────────────────────────────────────────
BU_RED      = RGBColor(0xCC, 0x00, 0x00)
DARK_GRAY   = RGBColor(0x2B, 0x2B, 0x2B)
MID_GRAY    = RGBColor(0x55, 0x55, 0x55)
LIGHT_GRAY  = RGBColor(0xF2, 0xF2, 0xF2)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GREEN       = RGBColor(0x1A, 0x7A, 0x3C)
ORANGE      = RGBColor(0xD6, 0x6B, 0x00)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Helpers ────────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def txb(slide, text, l, t, w, h,
        size=18, bold=False, color=DARK_GRAY, align=PP_ALIGN.LEFT,
        wrap=True, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def bullet_box(slide, items, l, t, w, h, size=16, color=DARK_GRAY, indent=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        prefix = "    • " if indent else "• "
        run.text = prefix + item
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb


def header_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 1.15, fill=BU_RED)
    txb(slide, title, 0.35, 0.12, 12, 0.55,
        size=32, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        txb(slide, subtitle, 0.35, 0.68, 12, 0.42,
            size=16, color=RGBColor(0xFF, 0xCC, 0xCC), align=PP_ALIGN.LEFT)
    add_rect(slide, 0, 1.15, 13.33, 6.35, fill=WHITE)


def divider(slide, y):
    add_rect(slide, 0.35, y, 12.63, 0.03, fill=LIGHT_GRAY)


# ── Slides ─────────────────────────────────────────────────────────────────────

def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 13.33, 7.5, fill=BU_RED)
    add_rect(slide, 0, 4.6, 13.33, 2.9, fill=DARK_GRAY)

    txb(slide, "NewtonBench", 0.6, 0.6, 12, 1.0,
        size=52, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    txb(slide, "Physics Law Discovery Benchmark for LLMs", 0.6, 1.6, 12, 0.7,
        size=24, color=RGBColor(0xFF, 0xCC, 0xCC), align=PP_ALIGN.LEFT)
    txb(slide, "Progress Update", 0.6, 2.4, 12, 0.6,
        size=20, italic=True, color=WHITE, align=PP_ALIGN.LEFT)

    txb(slide, "Karl Boma  •  Boston University", 0.6, 4.8, 8, 0.5,
        size=16, color=LIGHT_GRAY, align=PP_ALIGN.LEFT)
    txb(slide, "ExplorDSR  |  Science LLM Benchmarking", 0.6, 5.25, 8, 0.5,
        size=14, color=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.LEFT)


def slide_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, "What is NewtonBench?",
               "Benchmarking LLMs on scientific law discovery in simulated alien universes")

    txb(slide, "Core Idea", 0.5, 1.35, 12, 0.45,
        size=20, bold=True, color=BU_RED)
    txb(slide,
        "An LLM agent is placed in a simulated universe where the physical laws are deliberately "
        "altered from real-world physics. The agent must run experiments, analyze data, and submit "
        "a Python function encoding the discovered law — without using memorized formulas.",
        0.5, 1.8, 12.3, 0.9, size=15, color=DARK_GRAY)

    divider(slide, 2.8)

    col_items = [
        ("Modules (m0–m12)", [
            "One module per physics law",
            "Easy / Medium / Hard difficulty",
            "3 law versions (v0, v1, v2) per difficulty",
            "Alien constants & altered exponents",
        ]),
        ("Agent Backends", [
            "Vanilla agent — text only",
            "Code-assisted agent — Python execution",
            "Custom system prompt support",
        ]),
        ("Evaluation", [
            "RMSLE — numeric accuracy",
            "Exact Accuracy — correct functional form",
            "LLM symbolic equivalence judge",
        ]),
    ]

    x = 0.5
    for title, bullets in col_items:
        add_rect(slide, x, 2.95, 3.9, 3.9, fill=LIGHT_GRAY)
        txb(slide, title, x + 0.15, 3.05, 3.6, 0.45,
            size=15, bold=True, color=BU_RED)
        bullet_box(slide, bullets, x + 0.15, 3.55, 3.6, 3.0, size=13)
        x += 4.22


def slide_sourcer(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, "Wikipedia Law Sourcer",
               "Programmatically discovering new physics laws for the benchmark")

    txb(slide, "Motivation", 0.5, 1.35, 12, 0.4,
        size=18, bold=True, color=BU_RED)
    txb(slide,
        "NewtonBench currently has 12 hand-crafted modules. The goal: source new physics laws "
        "automatically from Wikipedia using an LLM extraction pipeline.",
        0.5, 1.75, 12.3, 0.6, size=14, color=DARK_GRAY)

    divider(slide, 2.5)

    txb(slide, "Pipeline  (utils/law_sourcer.py  +  source_laws.py)", 0.5, 2.6, 12, 0.4,
        size=18, bold=True, color=BU_RED)

    steps = [
        "Fetch ~193 law-candidate titles from Wikipedia's \"List of scientific laws\" article (MediaWiki API — no scraping, no new dependencies)",
        "For each candidate, fetch the plain-text article introduction",
        "Pass the intro to an LLM to extract: formula, variables, domain, and benchmark suitability (scalar output? 2–5 numeric inputs?)",
        "Filter to physics-only; flag laws already in NewtonBench",
        "Save structured JSON catalog",
    ]
    for i, s in enumerate(steps):
        add_rect(slide, 0.5, 3.1 + i * 0.72, 0.45, 0.45,
                 fill=BU_RED)
        txb(slide, str(i + 1), 0.5, 3.1 + i * 0.72, 0.45, 0.45,
            size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txb(slide, s, 1.1, 3.12 + i * 0.72, 11.7, 0.55,
            size=13, color=DARK_GRAY)


def slide_catalog(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, "Laws Catalog Results",
               "40 laws processed — 21 benchmarkable, 19 new to NewtonBench")

    # Stats boxes
    stats = [
        ("193", "Candidates\nfrom Wikipedia"),
        ("40", "Laws\nProcessed"),
        ("21", "Benchmarkable\nLaws"),
        ("19", "New to\nNewtonBench"),
    ]
    for i, (num, label) in enumerate(stats):
        x = 0.5 + i * 3.1
        add_rect(slide, x, 1.3, 2.8, 1.5, fill=BU_RED)
        txb(slide, num, x, 1.35, 2.8, 0.85,
            size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txb(slide, label, x, 2.15, 2.8, 0.6,
            size=12, color=RGBColor(0xFF, 0xCC, 0xCC), align=PP_ALIGN.CENTER)

    divider(slide, 3.05)

    txb(slide, "Sample New Benchmarkable Laws", 0.5, 3.15, 12, 0.4,
        size=17, bold=True, color=BU_RED)

    laws = [
        ("Archimedes' principle",   "fluid dynamics",     "F_b = ρ_f · g · V"),
        ("Beer–Lambert law",        "optics",             "A = ε · c · l"),
        ("Carnot efficiency",       "thermodynamics",     "η = (T_H − T_C) / T_H"),
        ("Hagen–Poiseuille",        "fluid dynamics",     "ΔP = 8μLQ / πR⁴"),
        ("Hubble's law",            "astrophysics",       "v = H₀ · D"),
        ("Doppler effect",          "waves & acoustics",  "f_obs = f_src · (v + v_obs) / (v + v_src)"),
    ]

    txb(slide, f"{'Law':<30}{'Domain':<22}Formula",
        0.5, 3.65, 12.3, 0.35, size=12, bold=True, color=MID_GRAY)
    add_rect(slide, 0.5, 3.98, 12.3, 0.03, fill=LIGHT_GRAY)

    for i, (name, domain, formula) in enumerate(laws):
        y = 4.05 + i * 0.48
        bg = LIGHT_GRAY if i % 2 == 0 else WHITE
        add_rect(slide, 0.5, y, 12.3, 0.46, fill=bg)
        txb(slide, name,    0.55, y + 0.05, 4.0,  0.38, size=12, color=DARK_GRAY)
        txb(slide, domain,  4.6,  y + 0.05, 3.5,  0.38, size=12, color=MID_GRAY, italic=True)
        txb(slide, formula, 8.1,  y + 0.05, 4.7,  0.38, size=12, color=BU_RED)


def slide_bernoulli(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, "New Module: m12_bernoulli",
               "Bernoulli's principle adapted for the alien-universe benchmark")

    txb(slide, "Design", 0.5, 1.35, 5.8, 0.4, size=18, bold=True, color=BU_RED)

    design = [
        "4 input variables: pressure, density, velocity, height",
        "Output: total flow energy quantity",
        "Hidden alien constant C₁ = 0.347 (real Bernoulli uses 0.5)",
        "3 difficulties × 3 versions = 9 distinct ground-truth laws",
    ]
    bullet_box(slide, design, 0.5, 1.8, 5.9, 2.0, size=13)

    txb(slide, "Law Variants", 6.7, 1.35, 6.0, 0.4, size=18, bold=True, color=BU_RED)

    variants = [
        ("Easy v0",   "B = C₁·ρ·v²  +  P",                           "no height term"),
        ("Easy v1",   "B = C₁·ρ·v¹·⁵  +  P",                        "fractional exponent"),
        ("Medium v0", "B = C₁·ρ·v²  +  P  +  C₂·ρ·h",               "adds height"),
        ("Medium v1", "B = C₁·ρ·v²·⁵  +  P  +  C₂·ρ·h",             "altered velocity exp"),
        ("Hard v0",   "B = C₁·ρ·v²  +  P  +  C₂·(ρ·h)¹·⁵",         "compound height"),
        ("Hard v1",   "B = C₁·ρ·v³  +  P¹·⁵  +  C₂·ρ·h²",          "higher order"),
    ]

    txb(slide, f"{'Variant':<14}{'Formula':<38}Note",
        6.7, 1.8, 6.1, 0.35, size=11, bold=True, color=MID_GRAY)
    add_rect(slide, 6.7, 2.12, 6.1, 0.03, fill=LIGHT_GRAY)

    for i, (v, f, n) in enumerate(variants):
        y = 2.18 + i * 0.52
        bg = LIGHT_GRAY if i % 2 == 0 else WHITE
        add_rect(slide, 6.7, y, 6.1, 0.5, fill=bg)
        txb(slide, v, 6.75, y + 0.05, 1.5, 0.42, size=11, bold=True, color=DARK_GRAY)
        txb(slide, f, 8.3,  y + 0.05, 3.0, 0.42, size=11, color=BU_RED)
        txb(slide, n, 11.3, y + 0.05, 1.4, 0.42, size=10, color=MID_GRAY, italic=True)

    divider(slide, 5.45)

    txb(slide, "Usage", 0.5, 5.55, 12, 0.4, size=18, bold=True, color=BU_RED)
    txb(slide,
        "python run_experiments.py --module m12_bernoulli --model_name ch45 "
        "--equation_difficulty easy --law_version v0 --trials 5 --agent_backend vanilla_agent",
        0.5, 6.0, 12.3, 0.55, size=12, color=MID_GRAY, italic=True)


def slide_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, "Prompt Comparison: Default vs. v9",
               "Haiku (ch45) on m12_bernoulli — easy difficulty, law v0, vanilla agent")

    # Table headers
    cols   = ["Trial", "Accuracy", "RMSLE", "Rounds", "Experiments", "Tokens"]
    widths = [0.65, 1.1, 1.05, 0.95, 1.5, 1.2]
    x0 = 0.4

    def table_section(title, color, rows, y_start):
        add_rect(slide, x0, y_start, 6.8, 0.38, fill=color)
        txb(slide, title, x0 + 0.1, y_start + 0.05, 6.6, 0.3,
            size=13, bold=True, color=WHITE)
        y = y_start + 0.38

        # header row
        x = x0
        for col, w in zip(cols, widths):
            add_rect(slide, x, y, w, 0.32, fill=DARK_GRAY)
            txb(slide, col, x + 0.04, y + 0.04, w - 0.08, 0.26,
                size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            x += w

        for ri, row in enumerate(rows):
            y += 0.32
            bg = LIGHT_GRAY if ri % 2 == 0 else WHITE
            x = x0
            for val, w in zip(row, widths):
                add_rect(slide, x, y, w, 0.32, fill=bg)
                c = GREEN if val == "YES" else (ORANGE if val == "NO" else DARK_GRAY)
                txb(slide, str(val), x + 0.04, y + 0.04, w - 0.08, 0.26,
                    size=10, bold=(val in ("YES", "NO")), color=c,
                    align=PP_ALIGN.CENTER)
                x += w

        # avg row
        y += 0.32
        x = x0
        for val, w in zip(rows[-1], widths):  # reuse last row slot
            add_rect(slide, x, y, w, 0.0, fill=WHITE)
            x += w
        return y

    default_rows = [
        ["0", "YES", "0.1716", "1", "0",  "7,123"],
        ["1", "YES", "0.1732", "1", "0",  "4,689"],
        ["2", "NO",  "0.1729", "7", "11", "49,117"],
        ["AVG", "67%", "0.1726", "3.0", "3.7", "—"],
    ]
    friend_rows = [
        ["0", "YES",  "0.0000", "6",  "37", "2,503"],
        ["1", "NO",   "0.3011", "6",  "14", "42,625"],
        ["2", "NO",   "1.3959", "2",  "3",  "10,569"],
        ["3", "NO",   "1.4056", "2",  "4",  "11,677"],
        ["4", "NO",   "1.4261", "2",  "3",  "10,298"],
        ["AVG", "20%", "0.9057", "3.6", "12.2", "—"],
    ]

    table_section("Default Prompt (BASE_PROMPT)", MID_GRAY, default_rows, 1.3)
    table_section("Friend's Prompt v9",           BU_RED,   friend_rows,  4.1)


def slide_findings(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, "Key Findings & Interpretation",
               "What the prompt comparison tells us")

    findings = [
        (GREEN,  "Default prompt: 67% accuracy — but it's not really discovery",
         [
             "Trials 0 & 1: 0 experiments, 1 round — Haiku pulled Bernoulli from training memory",
             "RMSLE ≈ 0.17 even on 'correct' trials: the alien constant (0.347) doesn't match the memorized one (0.5)",
             "The benchmark is being gamed, not solved",
         ]),
        (BU_RED, "Friend's v9 prompt: 20% accuracy — but trial 0 had RMSLE = 0.000",
         [
             "Trial 0: 37 experiments, 6 rounds — genuine experimental discovery, perfect result",
             "Trials 2–4: collapsed in 2 rounds with few experiments — Haiku struggles with the strict rules",
             "The prompt is better designed; the model isn't powerful enough to execute it consistently",
         ]),
        (ORANGE, "Conclusion",
         [
             "v9 prompt is the right direction — it prevents memorization and rewards real discovery",
             "Haiku is a weak model for this task; stronger models (Sonnet, GPT-4.1) will perform better",
             "v9 is now the default BASE_PROMPT in utils/vanilla_agent.py",
         ]),
    ]

    y = 1.35
    for color, title, bullets in findings:
        add_rect(slide, 0.4, y, 0.12, len(bullets) * 0.48 + 0.55, fill=color)
        txb(slide, title, 0.65, y + 0.05, 12.2, 0.42,
            size=15, bold=True, color=color)
        bullet_box(slide, bullets, 0.85, y + 0.5, 11.9, len(bullets) * 0.48,
                   size=13, color=DARK_GRAY)
        y += len(bullets) * 0.48 + 0.75


def slide_next_steps(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, "Next Steps", "Potential directions for expanding the benchmark")

    steps = [
        ("Expand the catalog",
         "Run the sourcer on the full 193 candidates and generate modules for the most "
         "interesting new laws (e.g. Hagen–Poiseuille, Carnot, Beer–Lambert)."),
        ("Auto-generate modules from catalog",
         "Build a second LLM step that takes a catalog entry and generates the full "
         "laws.py / core.py / prompts.py module code automatically."),
        ("Test stronger models with v9 prompt",
         "Run the same Bernoulli benchmark with Sonnet or GPT-4.1 to see if the "
         "genuine-discovery success rate improves."),
        ("Iterate prompt with optimizer",
         "Run run_iterative_optimization.py with the v9 prompt as the starting point "
         "to see if it can be further improved automatically."),
        ("Broader benchmark sweep",
         "Run all existing modules (m0–m12) with the new v9 prompt to get a full "
         "performance comparison against the old default."),
    ]

    for i, (title, desc) in enumerate(steps):
        y = 1.45 + i * 1.1
        add_rect(slide, 0.4, y, 0.55, 0.55, fill=BU_RED)
        txb(slide, str(i + 1), 0.4, y, 0.55, 0.55,
            size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txb(slide, title, 1.1, y,        12.0, 0.42, size=15, bold=True,  color=DARK_GRAY)
        txb(slide, desc,  1.1, y + 0.42, 12.0, 0.55, size=13, color=MID_GRAY)


# ── Build & save ───────────────────────────────────────────────────────────────

def build(output="NewtonBench_Progress.pptx"):
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_overview(prs)
    slide_sourcer(prs)
    slide_catalog(prs)
    slide_bernoulli(prs)
    slide_results(prs)
    slide_findings(prs)
    slide_next_steps(prs)

    prs.save(output)
    print(f"Saved: {output}")
    return output


if __name__ == "__main__":
    path = build()
    os.system(f'open "{path}"')
